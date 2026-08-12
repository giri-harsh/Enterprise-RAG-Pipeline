import logfire


def _parse_docx(file_path: str) -> str:
    """Paragraph text via python-docx."""
    import docx

    document = docx.Document(file_path)
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def _parse_pptx(file_path: str) -> str:
    """Shape text from every slide via python-pptx."""
    from pptx import Presentation

    texts = []
    for slide in Presentation(file_path).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def parse_office(file_path: str) -> str:
    """
    Extract text from .docx and .pptx.

    Uses python-docx and python-pptx directly, falling back to `unstructured`
    only for formats they do not cover.

    That ordering is deliberate. `unstructured` is the more capable library — it
    handles layout, tables and a much wider range of formats — but it is a very
    heavy dependency, and `evals/data_parser.py` already bypassed it for these two
    formats after hitting segfaults. Since the two narrow parsers were already in
    the tree and already trusted for exactly this job, using them on the main
    ingestion path too removes a large install from the default setup and makes
    the two code paths consistent.

    What this gives up: `unstructured` produces better-structured output for
    complex documents, particularly tables. For the plain prose in this corpus
    that difference does not show up in retrieval quality.
    """
    with logfire.span("Office document parsing", filename=file_path):
        extension = file_path.lower().rsplit(".", 1)[-1]

        try:
            if extension == "docx":
                text = _parse_docx(file_path)
            elif extension == "pptx":
                text = _parse_pptx(file_path)
            else:
                # Anything else — .doc, .odt, .ppt — needs the general partitioner.
                from unstructured.partition.auto import partition

                text = "\n".join(str(element) for element in partition(filename=file_path))

            if not text.strip():
                logfire.warning(f"No text extracted from {file_path}.")
            else:
                logfire.info(f"Extracted {len(text)} characters from {file_path}.")

            return text

        except ImportError as exc:
            raise RuntimeError(
                f"Cannot parse {file_path}: {exc}\n"
                "  pip install python-docx python-pptx\n"
                "  Other office formats also need: pip install unstructured"
            ) from exc

        except Exception as exc:
            logfire.error(f"Office parse failed for {file_path}: {exc}")
            raise
